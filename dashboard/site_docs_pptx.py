"""site-docs PPTX/PDF 생성기 — SiteDocs JSON → .pptx (+soffice PDF).

spec: docs/superpowers/specs/2026-06-24-site-docs-pptx-export-design.md
3 repo byte-identical: cop/hdel dashboard/, hermes scripts/.
"""
from __future__ import annotations

try:
    from pptx import Presentation as _Presentation  # noqa: F401
    from PIL import Image as _Image                  # noqa: F401
except ImportError as _e:                            # pragma: no cover
    import sys as _sys
    _sys.stderr.write(f"fatal: missing dependency ({_e}). Run: pip install python-pptx Pillow\n")
    raise SystemExit(2)

import json
import os
import re
import shutil
import subprocess
import sys

from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

_GRAY = RGBColor(0x6B, 0x72, 0x80)

_RUN_RE = re.compile(r"`([^`]+)`|\*\*([^*]+)\*\*|\[([^\]]+)\]\(([^)]+)\)")


def inline(text: str) -> list[dict]:
    """뷰어 render-markdown.ts inline() 미러: code / **bold** / [txt](url). 이탤릭 미지원.

    단일 레벨만(중첩 강조 없음). 링크는 표시텍스트만 유지하고 URL은 버림 — PPTX는
    비대화형이라 스킴 allowlist는 무의미(허용/비허용 모두 표시텍스트로 동일 렌더).
    """
    out: list[dict] = []
    pos = 0
    for m in _RUN_RE.finditer(text):
        if m.start() > pos:
            out.append({"t": "text", "s": text[pos:m.start()]})
        if m.group(1) is not None:                      # `code`
            out.append({"t": "code", "s": m.group(1)})
        elif m.group(2) is not None:                    # **bold**
            out.append({"t": "bold", "s": m.group(2)})
        else:                                           # [txt](url)
            out.append({"t": "link", "s": m.group(3)})  # display text only; url dropped
        pos = m.end()
    if pos < len(text):
        out.append({"t": "text", "s": text[pos:]})
    return out or [{"t": "text", "s": ""}]


_H_RE = re.compile(r"^(#{1,6})\s+(.+)$")
_UL_RE = re.compile(r"^[-*]\s+(.+)$")
_OL_RE = re.compile(r"^\d+\.\s+(.+)$")
_SEP_CELL_RE = re.compile(r"^:?-+:?$")


def _split_row(line: str) -> list[str]:
    s = line.strip()
    if s.startswith("|"):
        s = s[1:]
    if s.endswith("|"):
        s = s[:-1]
    return [c.strip() for c in s.split("|")]


def _is_table_sep(line: str) -> bool:
    cells = _split_row(line)
    return bool(cells) and all(_SEP_CELL_RE.match(c) for c in cells)


def _runs_text(runs: list[dict]) -> str:
    return "".join(r["s"] for r in runs)


def parse_markdown(md: str) -> list[dict]:
    if not md or not md.strip():
        return [{"type": "empty"}]
    blocks: list[dict] = []
    buf: list[str] = []
    cur_list = None  # ("ulist"|"olist", [items])
    in_code = False
    code_lines: list[str] = []

    def flush_para():
        nonlocal buf
        if buf:
            # join soft-wrapped lines, then strip — mirrors viewer flushBuf's buf.join(' ').trim()
            blocks.append({"type": "para", "runs": inline(" ".join(buf).strip())})
            buf = []

    def flush_list():
        nonlocal cur_list
        if cur_list:
            blocks.append({"type": cur_list[0], "items": cur_list[1]})
            cur_list = None

    lines = md.replace("\r\n", "\n").split("\n")
    i = 0
    while i < len(lines):
        raw = lines[i]
        line = raw.rstrip()
        if in_code:
            if line.strip() == "```":
                blocks.append({"type": "code", "text": "\n".join(code_lines)})
                code_lines = []
                in_code = False
            else:
                code_lines.append(raw)  # preserve leading whitespace; no inline
            i += 1
            continue
        if line.strip() == "```":
            flush_para(); flush_list(); in_code = True; i += 1; continue
        if line.lstrip().startswith("|") and i + 1 < len(lines) and _is_table_sep(lines[i + 1]):
            flush_para(); flush_list()
            header = [inline(c) for c in _split_row(line)]
            ncol = len(header)
            rows = []
            j = i + 2
            while (j < len(lines) and lines[j].lstrip().startswith("|")
                   and not _is_table_sep(lines[j])):     # a 2nd separator ends the table (viewer parity)
                cells = _split_row(lines[j])
                # normalize each row to header width: drop extras, pad missing (mirrors viewer)
                rows.append([inline(cells[c]) if c < len(cells) else [{"t": "text", "s": ""}]
                             for c in range(ncol)])
                j += 1
            blocks.append({"type": "table", "header": header, "rows": rows})
            i = j
            continue
        if not line.strip():
            flush_para(); flush_list(); i += 1; continue
        mh = _H_RE.match(line)
        if mh:
            flush_para(); flush_list()
            blocks.append({"type": "heading", "level": min(len(mh.group(1)), 3),
                           "runs": inline(mh.group(2))})
            i += 1; continue
        if line.startswith("> "):
            flush_para(); flush_list()
            blocks.append({"type": "quote", "runs": inline(line[2:])})
            i += 1; continue
        mu = _UL_RE.match(line)
        if mu:
            flush_para()
            if not cur_list or cur_list[0] != "ulist":
                flush_list(); cur_list = ("ulist", [])
            cur_list[1].append(inline(mu.group(1)))
            i += 1; continue
        mo = _OL_RE.match(line)
        if mo:
            flush_para()
            if not cur_list or cur_list[0] != "olist":
                flush_list(); cur_list = ("olist", [])
            cur_list[1].append(inline(mo.group(1)))
            i += 1; continue
        # paragraph buffer (soft-wrap)
        flush_list()
        buf.append(line)
        i += 1
    if in_code:                         # unterminated fence auto-closes
        blocks.append({"type": "code", "text": "\n".join(code_lines)})
    flush_para(); flush_list()
    return blocks or [{"type": "empty"}]


def pick_base_pt(total_chars: int) -> int:
    if total_chars < 500:
        return 12
    if total_chars < 1200:
        return 11
    if total_chars < 2500:
        return 10
    return 9


def slide_base_pt(slide_obj) -> int:
    return pick_base_pt(len(slide_obj.get("body_md", "")))


def resolve_screenshot(value: str, screenshots_dir: str) -> str | None:
    if not value or not screenshots_dir:
        return None
    name = os.path.basename(value)
    if "." not in name:                      # no extension -> not-found (no auto .png)
        return None
    path = os.path.join(screenshots_dir, name)
    return path if os.path.isfile(path) else None


def _has_decks(o) -> bool:
    return isinstance(o, dict) and ("spec" in o or "guide" in o)


def load_site_docs(raw) -> dict:
    """§3.3 우선순위: docs > top-level > 1겹 unwrap. 실패 시 ValueError."""
    if isinstance(raw, dict) and _has_decks(raw.get("docs")):
        return raw["docs"]
    if _has_decks(raw):
        return raw
    if isinstance(raw, dict):
        for k in ("data", "result"):
            if _has_decks(raw.get(k)):
                return raw[k]
    raise ValueError("SiteDocs not found: expected {spec,guide} or {docs:{...}} top-level")


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_run(r, run, base_pt, font, mono_font, *, bold=False, pt=None, color=None):
    r.text = run["s"] if isinstance(run, dict) else run
    is_code = isinstance(run, dict) and run["t"] == "code"
    is_bold = bold or (isinstance(run, dict) and run["t"] == "bold")
    r.font.name = mono_font if is_code else font
    r.font.size = Pt(pt if pt else base_pt)
    r.font.bold = is_bold
    if color is not None:
        r.font.color.rgb = color


def add_title_slide(prs, doc, opts):
    slide = _blank_slide(prs)
    n = len(doc.get("slides", []))
    title = (opts["title_prefix"] + " " if opts["title_prefix"] else "") + doc.get("title", "")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(12.333), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), {"t": "text", "s": title}, 32, opts["font"], opts["mono_font"], bold=True)
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(12.333), Inches(1.0))
    stf = sub.text_frame
    stf.word_wrap = True
    for i, txt in enumerate((_fmt_generated(doc.get("generatedAt", "")), f"슬라이드 {n}장")):
        para = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        _set_run(para.add_run(), {"t": "text", "s": txt}, 14, opts["font"], opts["mono_font"], color=_GRAY)


def _fmt_generated(iso: str) -> str:
    # 'YYYY-MM-DDTHH:MM:SS+09:00' -> 'YYYY-MM-DD HH:MM' (slice the local-stamped ISO)
    if len(iso) >= 16 and iso[10] == "T":
        return iso[:10] + " " + iso[11:16]
    return iso[:10] if iso else ""


# layout geometry (inches) — see spec §5.1
_MARGIN = 0.5
_BODY_TOP = 1.4
_BODY_H = 5.3
_FULL_W = 12.333


def layout_regions(layout):
    """(img|None, body) — inches dict {left, top, width, height}. 제목밴드/푸터 제외 영역."""
    top, h = _BODY_TOP, _BODY_H            # 1.4, 5.3
    if layout == "none":
        return None, {"left": _MARGIN, "top": top, "width": _FULL_W, "height": h}
    if layout == "top":
        img_h = h * 0.62
        gap = 0.12
        return ({"left": _MARGIN, "top": top, "width": _FULL_W, "height": img_h},
                {"left": _MARGIN, "top": top + img_h + gap, "width": _FULL_W, "height": h - img_h - gap})
    # side / split: 이미지 좌, 본문 우. split=균형, side=이미지 약간 좁게(세로 이미지)
    img_w = _FULL_W * (0.40 if layout == "side" else 0.48)
    gap = 0.3
    return ({"left": _MARGIN, "top": top, "width": img_w, "height": h},
            {"left": _MARGIN + img_w + gap, "top": top, "width": _FULL_W - img_w - gap, "height": h})


def _render_blocks(tf, blocks, base_pt, opts):
    """텍스트성 블록(heading/para/ulist/olist/quote/code/empty)을 한 텍스트 프레임에 단락으로."""
    first = True

    def para():
        nonlocal first
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        return p

    for b in blocks:
        if b["type"] == "heading":
            p = para()
            for run in b["runs"]:
                _set_run(p.add_run(), run, base_pt, opts["font"], opts["mono_font"],
                         bold=True, pt={1: 18, 2: 15, 3: 13}[b["level"]])
        elif b["type"] == "para":
            p = para()
            for run in b["runs"]:
                _set_run(p.add_run(), run, base_pt, opts["font"], opts["mono_font"])
        elif b["type"] in ("ulist", "olist"):
            for idx, item in enumerate(b["items"], 1):
                p = para()
                prefix = "• " if b["type"] == "ulist" else f"{idx}. "
                _set_run(p.add_run(), {"t": "text", "s": prefix}, base_pt, opts["font"], opts["mono_font"])
                for run in item:
                    _set_run(p.add_run(), run, base_pt, opts["font"], opts["mono_font"])
        elif b["type"] == "quote":
            p = para()
            p.level = 1
            for run in b["runs"]:
                _set_run(p.add_run(), run, base_pt, opts["font"], opts["mono_font"], color=_GRAY)
        elif b["type"] == "code":
            for ln in b["text"].split("\n"):
                p = para()
                _set_run(p.add_run(), {"t": "code", "s": ln}, base_pt, opts["font"], opts["mono_font"])
        elif b["type"] == "empty":
            p = para()
            _set_run(p.add_run(), {"t": "text", "s": "본문 없음"}, base_pt, opts["font"], opts["mono_font"], color=_GRAY)


def _add_title_band(slide, slide_obj, opts):
    tb = slide.shapes.add_textbox(Inches(_MARGIN), Inches(0.45), Inches(9.5), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    _set_run(p.add_run(), {"t": "text", "s": slide_obj.get("title", "")}, 24,
             opts["font"], opts["mono_font"], bold=True)
    cat = slide_obj.get("category", "")
    if cat:
        chip = slide.shapes.add_textbox(Inches(10.2), Inches(0.55), Inches(2.6), Inches(0.4))
        cp = chip.text_frame.paragraphs[0]
        cp.alignment = PP_ALIGN.RIGHT
        _set_run(cp.add_run(), {"t": "text", "s": cat}, 11, opts["font"], opts["mono_font"], color=_GRAY)


def _add_footer(slide, slide_obj, opts):
    date = _fmt_generated(slide_obj.get("updated_at", ""))[:10]
    commit = slide_obj.get("commit", "")
    text = f"updated {date}" + (f" · {commit}" if commit else "")
    fb = slide.shapes.add_textbox(Inches(_MARGIN), Inches(6.75), Inches(_FULL_W), Inches(0.3))
    p = fb.text_frame.paragraphs[0]
    _set_run(p.add_run(), {"t": "text", "s": text}, 9, opts["font"], opts["mono_font"], color=_GRAY)



def _slide_layout(slide_obj, png):
    lay = slide_obj.get("image_layout")
    if lay in ("top", "side", "split", "none"):
        return lay
    if not png:
        return "none"
    try:
        with Image.open(png) as im:
            r = im.size[0] / im.size[1]
    except Exception:
        return "split"
    return "top" if r >= 1.6 else ("side" if r <= 0.95 else "split")


def add_page_slide(prs, slide_obj, screenshots_dir, opts, progress_mode="none"):
    slide = _blank_slide(prs)
    _add_title_band(slide, slide_obj, opts)
    _add_footer(slide, slide_obj, opts)
    blocks = parse_markdown(slide_obj.get("body_md", ""))
    pt = slide_base_pt(slide_obj)
    png = resolve_screenshot(slide_obj.get("screenshot", ""), screenshots_dir) \
        if slide_obj.get("kind") != "system" else None
    layout = _slide_layout(slide_obj, png)
    if layout == "none":
        png = None
    img_region, body_region = layout_regions(layout)
    if png and img_region:
        try:
            _place_image(slide, png, img_region)
        except Exception:                       # corrupt/zero-byte/unsupported -> text-only + warn
            body_region = layout_regions("none")[1]
            emit(progress_mode, {"v": 1, "phase": "render", "status": "warn",
                                 "note": f"screenshot skipped: {slide_obj.get('id', '')}"})
    _render_body(slide, blocks, body_region["left"], body_region["width"], pt, opts,
                 top=body_region["top"], avail_h=body_region["height"])


def emit(progress_mode, event):
    if progress_mode == "none":
        return
    if progress_mode == "json":
        sys.stderr.write(json.dumps(event, ensure_ascii=False) + "\n")
        sys.stderr.flush()
        return
    # human
    phase = event.get("phase")
    if phase == "render":
        _dl = {"spec": "명세", "guide": "가이드"}.get(event.get("deck", ""), event.get("deck", ""))
        msg = f"[{_dl}] 슬라이드 {event.get('done')}/{event.get('total')} 렌더링…"
    elif phase == "pdf" and event.get("status") == "start":
        msg = "PDF 변환 중… (LibreOffice 첫 실행 5~15초)"
    elif phase == "done":
        warns = event.get("warnings") or []
        msg = "완료" + (f" (경고 {len(warns)}건)" if warns else "")
    else:
        msg = None
    if msg:
        sys.stderr.write(msg + "\n")
        sys.stderr.flush()


def _place_image(slide, png_path, region):
    with Image.open(png_path) as im:
        iw, ih = im.size
    box_w, box_h = Inches(region["width"]), Inches(region["height"])
    scale = min(box_w / iw, box_h / ih)   # scale 단위는 EMU/px (box_w=EMU, iw=px)
    w, h = int(iw * scale), int(ih * scale)   # EMU 단위 — 픽셀 아님
    left = Inches(region["left"]) + (box_w - w) // 2
    top = Inches(region["top"]) + (box_h - h) // 2
    pic = slide.shapes.add_picture(png_path, left, top, width=w, height=h)
    pic.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)   # thin border (spec §5.1)
    pic.line.width = Pt(0.75)


def _estimate_text_h(blocks, width_in, base_pt):
    # 줄높이·줄당 글자수를 base_pt에 맞춰 산정(한글 글리프는 ~1em 폭). 과소추정 시
    # 다음 블록(표 등)이 위 텍스트를 덮어써 폼이 깨짐 → 넉넉히 잡는다.
    line_h = max(0.20, base_pt / 72.0 * 1.5)
    cpl = max(8, int(width_in * (72.0 / max(base_pt, 1)) * 1.1))
    lines = 0.0
    for b in blocks:
        if b["type"] == "heading":
            lines += 1.5
        elif b["type"] in ("para", "quote"):
            lines += max(1, (len(_runs_text(b["runs"])) // cpl) + 1)
        elif b["type"] in ("ulist", "olist"):
            for it in b["items"]:
                lines += max(1, (len(_runs_text(it)) // cpl) + 1)
        elif b["type"] == "code":
            lines += b["text"].count("\n") + 1
        elif b["type"] == "empty":
            lines += 1
    return max(line_h, lines * line_h) + 0.1   # +0.1 텍스트박스 내부 여백


def _add_table(slide, table_block, x, y, w):
    header, rows = table_block["header"], table_block["rows"]
    ncol = max(1, max(len(header), *(len(r) for r in rows)) if rows else len(header))
    nrow = 1 + len(rows)
    gf = slide.shapes.add_table(nrow, ncol, Inches(x), Inches(y), Inches(w), Inches(0.4 * nrow))
    tbl = gf.table

    def fill(cell, runs):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        for run in runs:
            _set_run(p.add_run(), run, 11, "Apple SD Gothic Neo", "Menlo",
                     bold=(isinstance(run, dict) and run["t"] == "bold"))

    for c in range(ncol):
        fill(tbl.cell(0, c), header[c] if c < len(header) else [{"t": "text", "s": ""}])
    for r, row in enumerate(rows, 1):
        for c in range(ncol):
            fill(tbl.cell(r, c), row[c] if c < len(row) else [{"t": "text", "s": ""}])
    return Inches(0.4 * nrow)


def _estimate_total_h(blocks, w, base_pt):
    """본문 전체 추정 높이 — 세로 중앙 정렬용(_render_body 누적 로직 미러)."""
    total = 0.0
    group: list[dict] = []

    def flush():
        nonlocal total
        if group:
            total += max(0.3, _estimate_text_h(group, w, base_pt)) + 0.08
            group.clear()

    for b in blocks:
        if b["type"] == "table":
            flush()
            total += 0.4 * (1 + len(b["rows"])) + 0.1
        else:
            group.append(b)
    flush()
    return total


def _render_body(slide, blocks, x, w, base_pt, opts, top=_BODY_TOP, avail_h=_BODY_H):
    """플로우: 텍스트 블록은 텍스트 프레임에, 표는 별도 GraphicFrame. y 커서로 위→아래.
    짧은 본문은 세로 중앙 정렬(넘치면 상단부터)."""
    bottom = top + avail_h
    total = _estimate_total_h(blocks, w, base_pt)
    y = top + max(0.0, (avail_h - total) / 2)
    group: list[dict] = []

    def flush_group():
        nonlocal y
        if not group:
            return
        if bottom - y <= 0:
            group.clear(); return
        # 박스 높이로 y를 전진(추정값이 아니라 실제 그린 높이) — 안 그러면 다음 블록이 겹침.
        box_h = min(max(0.3, _estimate_text_h(group, w, base_pt)), bottom - y)
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(box_h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE   # bonus for interactive viewers
        _render_blocks(tf, list(group), base_pt, opts)
        y += box_h + 0.08   # 블록 간 간격
        group.clear()

    for b in blocks:
        if b["type"] == "table":
            flush_group()
            # always place the table — clip off-canvas if overflowing (recoverable in
            # PowerPoint), never silently drop it. clamp start so it stays near the slide.
            used = _add_table(slide, b, x, min(y, bottom), w)
            y += used / 914400 + 0.1
        else:
            group.append(b)
    flush_group()


def build_deck(doc, screenshots_dir, opts, progress_mode):
    slides = doc.get("slides", [])
    prs = _new_prs()
    add_title_slide(prs, doc, opts)
    for idx, s in enumerate(slides, 1):
        add_page_slide(prs, s, screenshots_dir, opts, progress_mode)
        emit(progress_mode, {"v": 1, "phase": "render", "deck": opts.get("_deck", ""),
                             "done": idx, "total": len(slides)})
    return prs


def convert_pdf(pptx_path, out_dir, progress_mode, soffice_bin="soffice"):
    deck = os.path.splitext(os.path.basename(pptx_path))[0]
    if not shutil.which(soffice_bin):
        emit(progress_mode, {"v": 1, "phase": "pdf", "deck": deck, "status": "warn",
                             "note": f"{soffice_bin} not found (install LibreOffice for PDF)"})
        return None
    emit(progress_mode, {"v": 1, "phase": "pdf", "deck": deck, "status": "start",
                         "note": "first run may take ~10s"})
    profile = f"/tmp/sd-soffice-{os.getpid()}"
    try:
        subprocess.run(
            [soffice_bin, "--headless", "--convert-to", "pdf", "--outdir", out_dir,
             f"-env:UserInstallation=file://{profile}", pptx_path],
            check=True, capture_output=True, timeout=120,
        )
        pdf = os.path.join(out_dir, deck + ".pdf")
        if os.path.isfile(pdf):
            emit(progress_mode, {"v": 1, "phase": "pdf", "deck": deck, "status": "ok"})
            return pdf
        emit(progress_mode, {"v": 1, "phase": "pdf", "deck": deck, "status": "warn",
                             "note": "soffice produced no pdf"})
        return None
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as e:
        emit(progress_mode, {"v": 1, "phase": "pdf", "deck": deck, "status": "warn",
                             "note": f"soffice failed: {type(e).__name__}"})
        return None
    finally:
        shutil.rmtree(profile, ignore_errors=True)


import argparse

_DECK_LABEL = {"spec": "기능명세", "guide": "사용가이드"}


def _read_raw(docs_arg):
    if docs_arg == "-":
        data = sys.stdin.read()
    else:
        with open(docs_arg, encoding="utf-8") as f:
            data = f.read()
    if not data.strip():
        raise ValueError("empty --docs input")
    return json.loads(data)


def main(argv=None):
    ap = argparse.ArgumentParser(description="site-docs PPTX/PDF generator")
    ap.add_argument("--proj", required=True)            # label only
    ap.add_argument("--docs", required=True)
    ap.add_argument("--screenshots-dir", default="")
    ap.add_argument("--out-dir", required=True)
    ap.add_argument("--deck", choices=["both", "spec", "guide"], default="both")
    ap.add_argument("--pdf", action="store_true")
    ap.add_argument("--progress", choices=["human", "json", "none"], default="human")
    ap.add_argument("--font", default="Apple SD Gothic Neo")
    ap.add_argument("--mono-font", default="Menlo")
    ap.add_argument("--title-prefix", default="")
    a = ap.parse_args(argv)
    pm = a.progress

    emit(pm, {"v": 1, "phase": "load", "status": "start"})
    try:
        site = load_site_docs(_read_raw(a.docs))
    except (ValueError, json.JSONDecodeError, OSError) as e:
        sys.stderr.write(f"fatal: {e}\n")
        return 2
    emit(pm, {"v": 1, "phase": "load", "status": "ok"})

    decks = ["spec", "guide"] if a.deck == "both" else [a.deck]
    opts = {"font": a.font, "mono_font": a.mono_font, "title_prefix": a.title_prefix, "label": a.proj}
    os.makedirs(a.out_dir, exist_ok=True)
    outputs, warnings = [], []
    if not a.screenshots_dir or not os.path.isdir(a.screenshots_dir):
        warnings.append("screenshots-dir absent: text-only decks")
        emit(pm, {"v": 1, "phase": "plan", "status": "warn", "note": "screenshots-dir absent: text-only"})
    emit(pm, {"v": 1, "phase": "plan", "decks": decks, "pdf": bool(a.pdf),
              "slideTotals": {k: len((site.get(k) or {}).get("slides", [])) for k in decks}})

    for deck in decks:
        doc = site.get(deck)
        if not doc:
            warnings.append(f"{deck}:missing deck")
            emit(pm, {"v": 1, "phase": "save", "deck": deck, "status": "warn", "note": "missing deck"})
            continue
        opts["_deck"] = deck
        prs = build_deck(doc, a.screenshots_dir, opts, pm)
        pptx_path = os.path.join(a.out_dir, f"{a.proj}-{_DECK_LABEL[deck]}.pptx")
        prs.save(pptx_path)
        outputs.append(pptx_path)
        emit(pm, {"v": 1, "phase": "save", "deck": deck, "status": "ok", "path": pptx_path})
        if a.pdf:
            pdf = convert_pdf(pptx_path, a.out_dir, pm)
            if pdf:
                outputs.append(pdf)
            else:
                warnings.append(f"pdf:{deck}:failed")

    emit(pm, {"v": 1, "phase": "done", "outputs": outputs, "warnings": warnings})
    for p in outputs:
        sys.stdout.write(p + "\n")
    sys.stdout.flush()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
