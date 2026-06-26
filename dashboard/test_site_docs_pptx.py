# test_site_docs_pptx.py
import io
import json
import os
import shutil
import subprocess
import tempfile
import unittest
from contextlib import redirect_stderr
import site_docs_pptx as sd
from pptx.util import Inches, Pt, Emu


class TestInline(unittest.TestCase):
    def test_plain_text(self):
        self.assertEqual(sd.inline("hello 한글"), [{"t": "text", "s": "hello 한글"}])

    def test_bold(self):
        self.assertEqual(
            sd.inline("a **b** c"),
            [{"t": "text", "s": "a "}, {"t": "bold", "s": "b"}, {"t": "text", "s": " c"}],
        )

    def test_code(self):
        self.assertEqual(
            sd.inline("run `x = 1` now"),
            [{"t": "text", "s": "run "}, {"t": "code", "s": "x = 1"}, {"t": "text", "s": " now"}],
        )

    def test_link_allowed_scheme_keeps_text_drops_url(self):
        self.assertEqual(
            sd.inline("see [docs](https://x.io/y)"),
            [{"t": "text", "s": "see "}, {"t": "link", "s": "docs"}],
        )

    def test_link_disallowed_scheme_stripped_to_text(self):
        # url group [^)]+ stops at the FIRST ')', so the trailing ')' is literal text.
        self.assertEqual(
            sd.inline("[bad](javascript:alert(1))"),
            [{"t": "link", "s": "bad"}, {"t": "text", "s": ")"}],
        )

    def test_italic_is_literal(self):
        self.assertEqual(sd.inline("a *b* c"), [{"t": "text", "s": "a *b* c"}])


class TestParseBlocks(unittest.TestCase):
    def test_heading_levels_clamped(self):
        b = sd.parse_markdown("# A\n## B\n### C\n#### D")
        self.assertEqual([(x["type"], x["level"]) for x in b],
                         [("heading", 1), ("heading", 2), ("heading", 3), ("heading", 3)])

    def test_paragraph_softwrap(self):
        b = sd.parse_markdown("line one\nline two\n\nnext para")
        self.assertEqual(b[0]["type"], "para")
        self.assertEqual(sd._runs_text(b[0]["runs"]), "line one line two")
        self.assertEqual(sd._runs_text(b[1]["runs"]), "next para")

    def test_unordered_list_flat(self):
        b = sd.parse_markdown("- a\n- b\n  - c")     # indented line is NOT nested
        ul = [x for x in b if x["type"] == "ulist"][0]
        self.assertEqual([sd._runs_text(i) for i in ul["items"]], ["a", "b"])
        self.assertTrue(any(x["type"] == "para" and sd._runs_text(x["runs"]) == "- c" for x in b))

    def test_ordered_list(self):
        b = sd.parse_markdown("1. first\n2. second")
        ol = [x for x in b if x["type"] == "olist"][0]
        self.assertEqual([sd._runs_text(i) for i in ol["items"]], ["first", "second"])

    def test_blockquote(self):
        b = sd.parse_markdown("> quoted")
        self.assertEqual(b[0]["type"], "quote")
        self.assertEqual(sd._runs_text(b[0]["runs"]), "quoted")

    def test_code_fence_raw_and_autoclose(self):
        b = sd.parse_markdown("```\n**not bold** 한글\nmore")
        code = [x for x in b if x["type"] == "code"][0]
        self.assertEqual(code["text"], "**not bold** 한글\nmore")

    def test_empty_body(self):
        self.assertEqual(sd.parse_markdown("   \n\n"), [{"type": "empty"}])


class TestParseTable(unittest.TestCase):
    def test_gfm_table(self):
        md = "| 이름 | 값 |\n| --- | --- |\n| a | **b** |"
        b = sd.parse_markdown(md)
        t = [x for x in b if x["type"] == "table"][0]
        self.assertEqual([sd._runs_text(c) for c in t["header"]], ["이름", "값"])
        self.assertEqual(sd._runs_text(t["rows"][0][0]), "a")
        self.assertEqual(t["rows"][0][1], [{"t": "bold", "s": "b"}])

    def test_pipe_without_separator_is_paragraph(self):
        b = sd.parse_markdown("| not | a table |\njust text")
        self.assertFalse(any(x["type"] == "table" for x in b))
        self.assertTrue(any(x["type"] == "para" for x in b))

    def test_table_then_text(self):
        md = "| h |\n| - |\n| x |\n\nafter"
        b = sd.parse_markdown(md)
        self.assertEqual([x["type"] for x in b], ["table", "para"])

    def test_table_ragged_rows_normalized_to_header(self):
        md = "| a | b |\n| - | - |\n| x | y | z |\n| only |"
        t = [x for x in sd.parse_markdown(md) if x["type"] == "table"][0]
        self.assertTrue(all(len(r) == 2 for r in t["rows"]))   # extra dropped, missing padded
        self.assertEqual(sd._runs_text(t["rows"][0][0]), "x")  # first cell kept
        self.assertEqual(sd._runs_text(t["rows"][1][1]), "")   # padded empty cell


class TestHelpers(unittest.TestCase):
    def test_pick_base_pt_bins(self):
        self.assertEqual(sd.pick_base_pt(100), 12)
        self.assertEqual(sd.pick_base_pt(800), 11)
        self.assertEqual(sd.pick_base_pt(2000), 10)
        self.assertEqual(sd.pick_base_pt(5000), 9)

    def test_resolve_screenshot_basename(self):
        with tempfile.TemporaryDirectory() as d:
            open(os.path.join(d, "home.png"), "wb").close()
            self.assertEqual(sd.resolve_screenshot("/static/cop/screenshots/home.png", d),
                             os.path.join(d, "home.png"))
            self.assertEqual(sd.resolve_screenshot("home.png", d), os.path.join(d, "home.png"))
            self.assertIsNone(sd.resolve_screenshot("missing.png", d))
            self.assertIsNone(sd.resolve_screenshot("", d))
            self.assertIsNone(sd.resolve_screenshot("noext", d))

    def test_load_site_docs_precedence(self):
        sg = {"spec": {"slides": []}, "guide": {"slides": []}}
        self.assertEqual(sd.load_site_docs({"docs": sg})["spec"], {"slides": []})   # docs wins
        self.assertEqual(sd.load_site_docs(sg)["guide"], {"slides": []})            # top-level
        self.assertEqual(sd.load_site_docs({"data": sg})["spec"], {"slides": []})   # one unwrap
        with self.assertRaises(ValueError):
            sd.load_site_docs({"nope": 1})


class TestDeckStructure(unittest.TestCase):
    OPTS = {"font": "Apple SD Gothic Neo", "mono_font": "Menlo", "title_prefix": "", "label": "cop"}

    def test_new_prs_is_16x9(self):
        prs = sd._new_prs()
        self.assertEqual(round(prs.slide_width / 914400, 3), 13.333)
        self.assertEqual(round(prs.slide_height / 914400, 3), 7.5)

    def test_title_slide_has_title_and_subtitle(self):
        prs = sd._new_prs()
        doc = {"title": "기능명세서 — CoP", "generatedAt": "2026-06-24T09:00:00+09:00",
               "slides": [{"id": "a"}, {"id": "b"}]}
        sd.add_title_slide(prs, doc, self.OPTS)
        texts = [sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame]
        self.assertTrue(any("기능명세서 — CoP" in t for t in texts))
        self.assertTrue(any("슬라이드 2장" in t for t in texts))   # N excludes title slide
        for sh in prs.slides[0].shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        self.assertEqual(r.font.name, "Apple SD Gothic Neo")


class TestPageSlideText(unittest.TestCase):
    OPTS = {"font": "Apple SD Gothic Neo", "mono_font": "Menlo", "title_prefix": "", "label": "cop"}

    def _slide(self, **over):
        s = {"id": "x", "kind": "system", "title": "아키텍처", "category": "시스템",
             "order": 1, "screenshot": "", "updated_at": "2026-06-24T09:00:00+09:00",
             "commit": "abc1234", "ui_hash": "", "body_md": "## 개요\n- 항목 하나\n- 항목 둘"}
        s.update(over); return s

    def test_system_slide_is_textonly_no_picture(self):
        prs = sd._new_prs()
        sd.add_page_slide(prs, self._slide(), "", self.OPTS)
        sl = prs.slides[0]
        self.assertFalse(any(sh.shape_type == 13 for sh in sl.shapes))  # 13 = PICTURE
        all_text = "\n".join(sh.text_frame.text for sh in sl.shapes if sh.has_text_frame)
        self.assertIn("아키텍처", all_text)
        self.assertIn("개요", all_text)
        self.assertIn("• 항목 하나", all_text)
        self.assertIn("abc1234", all_text)            # footer provenance

    def test_empty_body_placeholder(self):
        prs = sd._new_prs()
        sd.add_page_slide(prs, self._slide(body_md="  "), "", self.OPTS)
        all_text = "\n".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
        self.assertIn("본문 없음", all_text)


class TestTwoColumn(unittest.TestCase):
    OPTS = {"font": "Apple SD Gothic Neo", "mono_font": "Menlo", "title_prefix": "", "label": "cop"}

    def _png(self, d, w=640, h=480):
        from PIL import Image
        p = os.path.join(d, "home.png")
        Image.new("RGB", (w, h), (200, 200, 200)).save(p)
        return p

    def test_page_with_screenshot_has_picture_left_placed(self):
        with tempfile.TemporaryDirectory() as d:
            self._png(d)   # 640x480 (r=1.33) -> split layout: image on the LEFT, text right
            slide = {"id": "home", "kind": "page", "title": "홈", "category": "핵심",
                     "order": 1, "screenshot": "home.png", "image_layout": "split",
                     "updated_at": "2026-06-24T09:00:00+09:00",
                     "commit": "abc", "ui_hash": "", "body_md": "본문"}
            prs = sd._new_prs()
            sd.add_page_slide(prs, slide, d, self.OPTS)
            sl = prs.slides[0]
            pics = [sh for sh in sl.shapes if sh.shape_type == 13]
            self.assertEqual(len(pics), 1)
            self.assertAlmostEqual(pics[0].width / pics[0].height, 4 / 3, places=2)
            # adaptive split/side layout places the image in the LEFT region (not the old fixed right column)
            self.assertLess(pics[0].left + pics[0].width, prs.slide_width // 2 + Inches(1))

    def test_corrupt_png_falls_back_textonly(self):
        with tempfile.TemporaryDirectory() as d:
            open(os.path.join(d, "home.png"), "wb").close()    # 0-byte = corrupt
            slide = {"id": "home", "kind": "page", "title": "홈", "category": "핵심",
                     "order": 1, "screenshot": "home.png", "updated_at": "", "commit": "",
                     "ui_hash": "", "body_md": "본문"}
            prs = sd._new_prs()
            sd.add_page_slide(prs, slide, d, self.OPTS)         # must NOT raise
            pics = [sh for sh in prs.slides[0].shapes if sh.shape_type == 13]
            self.assertEqual(len(pics), 0)                      # no picture -> text-only


class TestNativeTable(unittest.TestCase):
    OPTS = {"font": "Apple SD Gothic Neo", "mono_font": "Menlo", "title_prefix": "", "label": "cop"}

    def test_table_block_becomes_native_table(self):
        slide = {"id": "t", "kind": "system", "title": "표", "category": "",
                 "order": 1, "screenshot": "", "updated_at": "2026-06-24T09:00:00+09:00",
                 "commit": "", "ui_hash": "",
                 "body_md": "설명 문단\n\n| 이름 | 값 |\n| --- | --- |\n| a | b |\n| c | d |"}
        prs = sd._new_prs()
        sd.add_page_slide(prs, slide, "", self.OPTS)
        tables = [sh for sh in prs.slides[0].shapes if sh.has_table]
        self.assertEqual(len(tables), 1)
        tbl = tables[0].table
        self.assertEqual(len(tbl.rows), 3)        # header + 2 body
        self.assertEqual(len(tbl.columns), 2)
        self.assertEqual(tbl.cell(0, 0).text, "이름")
        self.assertEqual(tbl.cell(2, 1).text, "d")

    def test_table_not_dropped_when_text_overflows(self):
        # heavy preceding text must not cause the trailing table to be silently dropped
        big = "\n\n".join(f"문단 {i} " + "내용 " * 20 for i in range(40))
        slide = {"id": "big", "kind": "system", "title": "큰 슬라이드", "category": "",
                 "order": 1, "screenshot": "", "updated_at": "", "commit": "", "ui_hash": "",
                 "body_md": big + "\n\n| 키 | 값 |\n| - | - |\n| a | b |"}
        prs = sd._new_prs()
        sd.add_page_slide(prs, slide, "", self.OPTS)
        tables = [sh for sh in prs.slides[0].shapes if sh.has_table]
        self.assertEqual(len(tables), 1)   # present even though text overflowed (clip, not drop)


class TestBuildDeck(unittest.TestCase):
    OPTS = {"font": "Apple SD Gothic Neo", "mono_font": "Menlo", "title_prefix": "", "label": "cop"}

    def test_deck_slide_count(self):
        doc = {"title": "기능명세서 — CoP", "generatedAt": "2026-06-24T09:00:00+09:00",
               "slides": [
                   {"id": "a", "kind": "page", "title": "A", "category": "핵심", "order": 1,
                    "screenshot": "", "updated_at": "", "commit": "", "ui_hash": "", "body_md": "본문 A"},
                   {"id": "b", "kind": "system", "title": "B", "category": "시스템", "order": 2,
                    "screenshot": "", "updated_at": "", "commit": "", "ui_hash": "", "body_md": "본문 B"},
               ]}
        prs = sd.build_deck(doc, "", self.OPTS, "none")
        self.assertEqual(len(prs.slides), 3)       # title + 2


class TestProgress(unittest.TestCase):
    def test_json_event_is_ndjson_with_version(self):
        buf = io.StringIO()
        with redirect_stderr(buf):
            sd.emit("json", {"v": 1, "phase": "render", "deck": "spec", "done": 1, "total": 3})
        line = buf.getvalue().strip()
        obj = json.loads(line)
        self.assertEqual(obj["v"], 1)
        self.assertEqual(obj["phase"], "render")

    def test_none_is_silent(self):
        buf = io.StringIO()
        with redirect_stderr(buf):
            sd.emit("none", {"v": 1, "phase": "done", "outputs": []})
        self.assertEqual(buf.getvalue(), "")


class TestPdf(unittest.TestCase):
    @unittest.skipUnless(shutil.which("soffice"), "soffice not installed")
    def test_convert_pdf_creates_nonempty_pdf(self):
        with tempfile.TemporaryDirectory() as d:
            prs = sd._new_prs()
            sd.add_title_slide(prs, {"title": "T", "generatedAt": "", "slides": []},
                               {"font": "Apple SD Gothic Neo", "mono_font": "Menlo",
                                "title_prefix": "", "label": "t"})
            pptx_path = os.path.join(d, "t.pptx")
            prs.save(pptx_path)
            pdf = sd.convert_pdf(pptx_path, d, "none")
            self.assertTrue(pdf and os.path.getsize(pdf) > 0)
            self.assertTrue(pdf.endswith(".pdf"))

    def test_convert_pdf_missing_soffice_returns_none(self):
        with tempfile.TemporaryDirectory() as d:
            open(os.path.join(d, "x.pptx"), "wb").close()
            self.assertIsNone(sd.convert_pdf(os.path.join(d, "x.pptx"), d, "none",
                                             soffice_bin="definitely-not-soffice"))


class TestMain(unittest.TestCase):
    def _docs_json(self):
        d = {"spec": {"title": "기능명세서 — CoP", "generatedAt": "2026-06-24T09:00:00+09:00",
                      "slides": [{"id": "a", "kind": "system", "title": "A", "category": "",
                                  "order": 1, "screenshot": "", "updated_at": "", "commit": "",
                                  "ui_hash": "", "body_md": "본문"}]},
             "guide": {"title": "사용가이드 — CoP", "generatedAt": "2026-06-24T09:00:00+09:00",
                       "slides": [{"id": "a", "kind": "system", "title": "A", "category": "",
                                   "order": 1, "screenshot": "", "updated_at": "", "commit": "",
                                   "ui_hash": "", "body_md": "본문"}]}}
        return d

    def test_main_writes_two_decks_and_prints_paths(self):
        with tempfile.TemporaryDirectory() as d:
            jp = os.path.join(d, "data.json")
            with open(jp, "w", encoding="utf-8") as f:
                json.dump({"docs": self._docs_json()}, f)
            out = os.path.join(d, "exports")
            buf = io.StringIO()
            from contextlib import redirect_stdout
            with redirect_stdout(buf):
                rc = sd.main(["--proj", "cop", "--docs", jp, "--screenshots-dir", d,
                              "--out-dir", out, "--progress", "none"])
            self.assertEqual(rc, 0)
            paths = buf.getvalue().split()
            self.assertEqual(len([p for p in paths if p.endswith(".pptx")]), 2)
            self.assertTrue(os.path.isfile(os.path.join(out, "cop-기능명세.pptx")))
            self.assertTrue(os.path.isfile(os.path.join(out, "cop-사용가이드.pptx")))

    def test_main_bad_json_exits_2(self):
        with tempfile.TemporaryDirectory() as d:
            jp = os.path.join(d, "bad.json")
            with open(jp, "w") as f:
                f.write('{"nope": 1}')
            rc = sd.main(["--proj", "cop", "--docs", jp, "--screenshots-dir", d,
                          "--out-dir", d, "--progress", "none"])
            self.assertEqual(rc, 2)


class TestIntegrationCop(unittest.TestCase):
    DATA = "data.json"
    SHOTS = "screenshots"

    @unittest.skipUnless(os.path.isfile("data.json"), "run from cop/dashboard with built data.json")
    def test_real_cop_build_with_pdf(self):
        with tempfile.TemporaryDirectory() as out:
            rc = sd.main(["--proj", "cop", "--docs", self.DATA, "--screenshots-dir", self.SHOTS,
                          "--out-dir", out, "--pdf", "--progress", "none"])
            self.assertEqual(rc, 0)
            spec_pptx = os.path.join(out, "cop-기능명세.pptx")
            self.assertTrue(os.path.isfile(spec_pptx))
            from pptx import Presentation
            prs = Presentation(spec_pptx)
            self.assertTrue(any(sh.shape_type == 13 for sl in prs.slides for sh in sl.shapes))
            spec_pdf = os.path.join(out, "cop-기능명세.pdf")
            if shutil.which("soffice") and shutil.which("pdftotext") and os.path.isfile(spec_pdf):
                txt = subprocess.run(["pdftotext", spec_pdf, "-"], capture_output=True,
                                     text=True).stdout
                self.assertRegex(txt, r"[가-힣]")     # real Hangul, not tofu/blank


class TestAdaptiveLayout(unittest.TestCase):
    """슬라이드 재설계: image_layout 자동분류 + per-slide base_pt + 영역 계산."""

    def test_slide_base_pt_per_slide(self):
        self.assertEqual(sd.slide_base_pt({"body_md": "x" * 100}), 12)
        self.assertEqual(sd.slide_base_pt({"body_md": "x" * 3000}), 9)

    def test_layout_regions_top_image_above_body(self):
        img, body = sd.layout_regions("top")
        self.assertLess(img["top"], body["top"])

    def test_layout_regions_side_image_left_of_body(self):
        img, body = sd.layout_regions("side")
        self.assertLess(img["left"], body["left"])

    def test_layout_regions_split_wider_than_side(self):
        # split=균형(48%)이 side=세로용(40%)보다 이미지 폭 넓음
        self.assertGreater(sd.layout_regions("split")[0]["width"],
                           sd.layout_regions("side")[0]["width"])

    def test_layout_regions_none_has_no_image(self):
        img, body = sd.layout_regions("none")
        self.assertIsNone(img)

    def test_slide_layout_explicit_field_wins(self):
        self.assertEqual(sd._slide_layout({"image_layout": "side"}, None), "side")

    def test_slide_layout_no_png_is_none(self):
        self.assertEqual(sd._slide_layout({}, None), "none")

    def test_slide_layout_pil_fallback_by_ratio(self):
        from PIL import Image
        with tempfile.TemporaryDirectory() as d:
            wide = os.path.join(d, "wide.png"); Image.new("RGB", (1280, 720)).save(wide)   # 1.78 -> top
            tall = os.path.join(d, "tall.png"); Image.new("RGB", (418, 628)).save(tall)    # 0.67 -> side
            sq = os.path.join(d, "sq.png"); Image.new("RGB", (760, 594)).save(sq)          # 1.28 -> split
            self.assertEqual(sd._slide_layout({}, wide), "top")
            self.assertEqual(sd._slide_layout({}, tall), "side")
            self.assertEqual(sd._slide_layout({}, sq), "split")


if __name__ == "__main__":
    unittest.main()
