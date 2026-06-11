"""월별 활동보고서 pptx 생성 (frontmatter → 5월 템플릿 텍스트 슬롯 교체).

원본 디자인 100% 유지. 슬라이드/Shape 추가/삭제 없이 텍스트 run 만 교체.
사용:
  .venv/bin/python3 scripts/build_report_pptx.py --month 2026-06 [--out path.pptx]

요구:
  - assets/reports/templates/2026-05_template.pptx (마스터)
  - Obsidian 의 CoP_PhysicalAI_{month}_활동보고서.md (frontmatter 데이터)
"""
import argparse
import copy
import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Emu

REPO_ROOT = Path(__file__).resolve().parent.parent
TEMPLATE_PATH = REPO_ROOT / "assets/reports/templates/2026-05_template.pptx"
OBSIDIAN_DIR = (
    Path.home() / "Documents/second-brain/03 Areas/회사문서/CoP_PhysicalAI"
)


# ---------- frontmatter 로더 ----------
def load_frontmatter(month: str) -> dict:
    """Obsidian 보고서 .md 의 YAML frontmatter 로드."""
    md = OBSIDIAN_DIR / f"CoP_PhysicalAI_{month}_활동보고서.md"
    if not md.exists():
        raise FileNotFoundError(f"보고서 없음: {md}")
    text = md.read_text(encoding="utf-8")
    m = re.match(r"^---\s*\n(.*?)\n---\s*\n", text, re.DOTALL)
    if not m:
        raise ValueError(f"frontmatter 없음: {md}")
    return yaml.safe_load(m.group(1))


# ---------- 텍스트 교체 ----------
def _patch_text_frame(tf, new_text: str):
    """text_frame 내 첫 run 스타일 유지하면서 텍스트 교체."""
    lines = (new_text or "").split("\n")
    # 모든 paragraph 들에서 첫 paragraph 만 남김
    paragraphs = list(tf.paragraphs)
    if not paragraphs:
        return
    first_p = paragraphs[0]
    for p in paragraphs[1:]:
        p._p.getparent().remove(p._p)
    # 첫 paragraph 의 첫 run 만 남기고 나머지 제거 (스타일은 첫 run 이 보존)
    runs = list(first_p.runs)
    if runs:
        first_run = runs[0]
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
        first_run.text = lines[0] if lines else ""
    else:
        first_p.text = lines[0] if lines else ""
    # 추가 줄: 첫 paragraph 를 복제해서 텍스트만 교체
    for extra in lines[1:]:
        new_p = copy.deepcopy(first_p._p)
        first_p._p.getparent().append(new_p)
        # 새 paragraph 의 첫 run.text 만 교체
        # lxml elem 직접 접근
        from pptx.oxml.ns import qn
        rs = new_p.findall(qn("a:r"))
        if rs:
            t = rs[0].find(qn("a:t"))
            if t is not None:
                t.text = extra
            # 나머지 run 제거
            for r in rs[1:]:
                new_p.remove(r)
        # 결과 paragraph 가 tf 의 맨 끝이 되도록 했음 (append) — OK


def replace_text(shape, new_text: str):
    if not shape.has_text_frame:
        return
    _patch_text_frame(shape.text_frame, new_text)


def replace_cell(cell, new_text: str):
    """표 셀 텍스트 교체."""
    _patch_text_frame(cell.text_frame, new_text)


def find_by_name(slide, name: str):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


# ---------- 데이터 포맷 헬퍼 ----------
def period_dashed(period: str) -> str:
    """'2026.06.01 ~ 2026.06.30' 같은 그대로 반환."""
    return period or ""


def month_num(month: str) -> str:
    """'2026-06' → '6'."""
    return str(int(month.split("-")[1])) if month else ""


def fmt_achievement(a: dict) -> str:
    title = a.get("title", "")
    value = a.get("value", "")
    return f"{title}: {value}" if title and value else (title or value or "")


# ---------- 슬라이드별 매핑 ----------
def patch_slide1(slide, fm: dict):
    """표지 슬라이드."""
    mn = month_num(fm["month"])
    status = fm.get("status", {})
    ovw = fm.get("overview", {})
    ach = fm.get("achievements", [])

    mapping = {
        "Text 4": f"2026년 {mn}월",
        "Text 5": f"보고 기간: {period_dashed(fm.get('period',''))}",
        "Text 6": f"2026년 사내 CoP 활동보고서 — {mn}월",
        "Text 7": fm.get("subtitle", ""),
        "Text 11": fm.get("status_label", ""),
        "Text 12": f"CoP 명칭: {ovw.get('cop_name','')}",
        "Text 13": f"보고 회차: {fm.get('report_round','')}",
        "Text 14": f"활동 주제: {ovw.get('activity_topic','')}",
        "Text 15": f"최종 목표: {ovw.get('final_goal','')}",
        "Text 18": status.get("state", ""),
        "Text 19": f"진행률: {status.get('progress_pct',0)}% ({status.get('progress_weeks','')})",
        "Text 20": status.get("detail", ""),
    }
    for i in range(3):
        mapping[f"Text {23+i}"] = fmt_achievement(ach[i]) if i < len(ach) else ""
    for name, val in mapping.items():
        sh = find_by_name(slide, name)
        if sh:
            replace_text(sh, val)


def patch_slide2(slide, fm: dict):
    """CoP 개요 슬라이드 (표 + 활동현황)."""
    mn = month_num(fm["month"])
    status = fm.get("status", {})
    ovw = fm.get("overview", {})

    mapping = {
        "Text 7": f"2026년 {mn}월",
        "Text 8": f"보고 기간: {period_dashed(fm.get('period',''))}",
        "Text 9": fm.get("status_label", ""),
        "Text 15": fm.get("status_label", ""),
        "Text 17": status.get("state", ""),
        "Text 18": f"{status.get('progress_pct',0)}%",
        "Text 19": status.get("detail", ""),
    }
    for name, val in mapping.items():
        sh = find_by_name(slide, name)
        if sh:
            replace_text(sh, val)

    # 표 (6x2) — R1~R5 데이터 교체
    table_shape = next((s for s in slide.shapes if s.has_table), None)
    if table_shape:
        rows = list(table_shape.table.rows)
        cells_data = [
            ("CoP 명칭", ovw.get("cop_name", "")),
            ("보고 기간", period_dashed(fm.get("period", ""))),
            ("보고 회차", fm.get("report_round", "")),
            ("활동 주제", ovw.get("activity_topic", "")),
            ("최종 목표", ovw.get("final_goal", "")),
        ]
        for i, (label, val) in enumerate(cells_data, start=1):
            if i < len(rows):
                cells = list(rows[i].cells)
                if len(cells) >= 2:
                    replace_cell(cells[0], label)
                    replace_cell(cells[1], val)


def patch_slide3(slide, fm: dict):
    """활동 상세 — 활동요약 + 주차별."""
    mn = month_num(fm["month"])
    weeks = fm.get("weeks", []) or []

    mapping = {
        "Text 3": f"{mn}월 활동 및 성과",
        "Text 6": fm.get("activity_summary", ""),
        "Text 9": fm.get("overview", {}).get("final_goal", ""),
    }
    # 주차 슬롯 (week / date / content) — 5월 템플릿 기준 5주차
    # Text 22, 24, 27 — 1주차
    # Text 29, 31, 34 — 2주차
    # Text 36, 38, 41 — 3주차
    # Text 43, 45, 48 — 4주차
    # Text 50, 52, 55 — 5주차
    week_slots = [
        ("Text 22", "Text 24", "Text 27"),
        ("Text 29", "Text 31", "Text 34"),
        ("Text 36", "Text 38", "Text 41"),
        ("Text 43", "Text 45", "Text 48"),
        ("Text 50", "Text 52", "Text 55"),
    ]
    for i, (wn, wd, wc) in enumerate(week_slots):
        if i < len(weeks):
            w = weeks[i]
            mapping[wn] = str(w.get("week", ""))
            mapping[wd] = str(w.get("date", ""))
            mapping[wc] = str(w.get("content", ""))
        else:
            mapping[wn] = ""
            mapping[wd] = ""
            mapping[wc] = ""
    for name, val in mapping.items():
        sh = find_by_name(slide, name)
        if sh:
            replace_text(sh, val)


def patch_slide4(slide, fm: dict):
    """차월 계획 — 카테고리 2개, 각 카테고리 3 items."""
    nm = fm.get("next_month", {}) or {}
    mapping = {
        "Text 3": nm.get("title", ""),
        "Text 4": nm.get("title", ""),
        "Text 7": f"계획 기간: {nm.get('period','')}",
    }

    cats = nm.get("categories", []) or []
    # 카테고리 1 (5월 템플릿 기준):
    # Text 10 = label, Text 12 = state
    # 3 items: (Text 15, 16) (Text 19, 20) (Text 23, 24)
    # 카테고리 2:
    # Text 27 = label, Text 29 = state
    # 3 items: (Text 32, 33) (Text 36, 37) (Text 40, 41)
    cat_slots = [
        ("Text 10", "Text 12", [("Text 15", "Text 16"), ("Text 19", "Text 20"), ("Text 23", "Text 24")]),
        ("Text 27", "Text 29", [("Text 32", "Text 33"), ("Text 36", "Text 37"), ("Text 40", "Text 41")]),
    ]
    for ci, (label_slot, state_slot, item_slots) in enumerate(cat_slots):
        if ci < len(cats):
            c = cats[ci]
            mapping[label_slot] = c.get("label", "")
            mapping[state_slot] = c.get("state", "")
            items = c.get("items", []) or []
            for ii, (t_slot, d_slot) in enumerate(item_slots):
                if ii < len(items):
                    it = items[ii]
                    mapping[t_slot] = it.get("title", "")
                    mapping[d_slot] = it.get("desc", "")
                else:
                    mapping[t_slot] = ""
                    mapping[d_slot] = ""
        else:
            mapping[label_slot] = ""
            mapping[state_slot] = ""
            for t_slot, d_slot in item_slots:
                mapping[t_slot] = ""
                mapping[d_slot] = ""
    for name, val in mapping.items():
        sh = find_by_name(slide, name)
        if sh:
            replace_text(sh, val)


def patch_slide5(slide, fm: dict):
    """활동 현장 사진. site_caption 만 동적 교체. 사진 자체는 frontmatter site_image 가 있으면 교체."""
    cap = fm.get("site_caption", "")
    sh = find_by_name(slide, "Text 6")
    if sh and cap:
        replace_text(sh, cap)

    # 사진 교체: site_image 가 로컬 파일이면 교체. /static/cop/... 같은 URL 은 스킵.
    site_img = fm.get("site_image", "")
    if site_img and not site_img.startswith(("http", "/static")):
        img_path = REPO_ROOT / site_img.lstrip("/")
        if img_path.exists():
            replace_picture(slide, img_path)


def replace_picture(slide, new_img_path: Path):
    """첫 번째 큰 사진 (활동 현장) 을 새 이미지로 교체."""
    # 가장 큰 picture shape 를 찾아 교체 (또는 Image 0 같은 이름)
    pics = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
    if not pics:
        return
    target = max(pics, key=lambda p: (p.width or 0) * (p.height or 0))
    # 위치/크기 유지하면서 교체
    left, top, w, h = target.left, target.top, target.width, target.height
    sp = target._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(str(new_img_path), left, top, width=w, height=h)


# ---------- 메인 ----------
def build(month: str, out_path: Path):
    fm = load_frontmatter(month)
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"템플릿 없음: {TEMPLATE_PATH}")
    p = Presentation(str(TEMPLATE_PATH))
    if len(p.slides) < 5:
        raise ValueError(f"템플릿 슬라이드 부족: {len(p.slides)} (5 필요)")
    patch_slide1(p.slides[0], fm)
    patch_slide2(p.slides[1], fm)
    patch_slide3(p.slides[2], fm)
    patch_slide4(p.slides[3], fm)
    patch_slide5(p.slides[4], fm)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(out_path))
    return out_path


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--month", required=True, help="YYYY-MM (예: 2026-06)")
    ap.add_argument("--out", default=None, help="출력 경로 (기본: dist/2026_cop_physical_AI_{month}.pptx)")
    args = ap.parse_args()

    out = Path(args.out) if args.out else REPO_ROOT / f"dist/2026_cop_physical_AI_{args.month}.pptx"
    try:
        result = build(args.month, out)
        print(f"✓ pptx 생성 완료: {result} ({result.stat().st_size:,} bytes)")
    except Exception as e:
        print(f"✗ 실패: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
